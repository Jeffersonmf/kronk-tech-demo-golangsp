"""Builds the GolangSP "IA sem Nuvem, em Golang" deck as a highly professional,
dense, and visually stunning widescreen .pptx presentation.

This script features a premium dark-mode theme, custom container layouts,
VS Code-like monospace terminal windows for code blocks, alternating row
dashboard styling for tables, automated key-concept highlighting for bullets,
dynamic section divider slides (chapters), and elegant headers/footers with page numbers.

Source of truth for content is the hardcoded SLIDES array, which corresponds to:
slides-abertura.md, slides-kronk-deepdive.md, slides-demos-finais.md.
Full spoken scripts go into PPTX speaker notes.

Run: python3 slides/build_deck.py
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

OUT_PATH = "slides/GolangSP-Kronk-Talk.pptx"
GOPHER_PATH = "slides/assets/gopher_new.png"

# ==============================================================================
# VIBRANT TECH BLUE & CRISP WHITE THEME (HIGH CONTRAST, ZERO GLOOMINESS)
# ==============================================================================
BG_COLOR = RGBColor(0x0A, 0x25, 0x40)        # Stripe-style Deep Tech Blue background
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)         # Crisp Pure White container cards
CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)     # Subtle Light Gray card borders (Slate-200)
CODE_BG = RGBColor(0x0F, 0x17, 0x2A)         # Dark Slate for realistic code terminal background
CODE_BORDER = RGBColor(0x33, 0x41, 0x55)     # Slate-700 terminal border
TEXT_MAIN = RGBColor(0x0F, 0x17, 0x2A)       # Deep Charcoal Slate-900 for pristine legibility on white cards
TEXT_MUTED = RGBColor(0x47, 0x55, 0x69)      # Slate-600 secondary text on white cards
ACCENT_BLUE = RGBColor(0x00, 0x6F, 0xE6)     # Brilliant Royal Go Blue for active highlights
ACCENT_ORANGE = RGBColor(0xEF, 0x44, 0x44)   # Red (Rose-500) for contrast warnings

# ==============================================================================
# SECTION DIVIDER METADATA
# ==============================================================================
SECTIONS_META = {
    "intro": {
        "num": "01",
        "title": "CONCEITO & CONTEXTO",
        "sub": "O desafio dos custos na nuvem e a revolução da IA local-first"
    },
    "kronk": {
        "num": "02",
        "title": "DEEP DIVE: KRONK SDK",
        "sub": "Orquestrando quantização, cache de contexto, MoE e hardware híbrido em Go"
    },
    "code": {
        "num": "03",
        "title": "POR DENTRO DO CÓDIGO",
        "sub": "Análise técnica comparativa: do servidor MCP ingênuo ao indexador semântico"
    },
    "demo": {
        "num": "04",
        "title": "A PROVA REAL: DEMOS",
        "sub": "Cena 1 (nuvem + dumb) contra Cena 2 (local + smart) ao vivo na tela"
    },
    "outro": {
        "num": "05",
        "title": "CONCLUINDO A JORNADA",
        "sub": "Key takeaways, caminhos de escala e o futuro do desenvolvedor Go"
    }
}

# ==============================================================================
# RAW SOURCE OF TRUTH: SLIDES DATA (including speaker notes, tables, code)
# ==============================================================================
SLIDES = [
    {
        "kind": "title",
        "title": "AI sem Cloud com Golang!",
        "subtitle": "Golang + AI + MCP — Evento Oficial GolangSP",
        "notes": "Boa noite, pessoal! É uma honra imensa estar aqui no evento oficial da GolangSP. "
                 "Hoje vamos falar de um tema fascinante e de altíssimo nível: "
                 "como criar e rodar inteligência artificial sem depender de serviços em nuvem, usando Go!"
    },
    {
        "kind": "bio",
        "title": "Jefferson Ferreira",
        "role": "Tech Lead & Staff Software Engineer",
        "companies": [
            {"name": "1Global (Atual)", "desc": "Senior Software Engineer — Liderando arquitetura, escalabilidade e infraestrutura de sistemas modernos."},
            {"name": "Trampoline / Tino / Tractian", "desc": "Staff & Principal Engineer — Arquitetura de IA (RAG/MCP), alto throughput e pipelines distribuídos em Go/Rust."},
            {"name": "IBM / XP Inc. / Truphone", "desc": "Senior Systems Architect — Resiliência, Big Data, infraestrutura de nuvem e engenharia enterprise."},
            {"name": "B3 / Itaú / Xerox", "desc": "Systems & Security Architect — Sistemas financeiros pós-trade de missão crítica, mobile banking e criptografia."}
        ],
        "notes": "Antes de mergulharmos de cabeça no código e na infraestrutura técnica, deixa eu me apresentar rapidamente. "
                 "Sou o Jefferson, atuo como Engenheiro de Software Sênior e Tech Lead. "
                 "Ao longo da minha trajetória profissional de mais de 20 anos, liderei arquitetura e engenharia em frentes de alta complexidade: "
                 "atualmente estou na 1Global, focando em sistemas escaláveis de alto desempenho em Go; "
                 "como consultor independente e em empresas como Tino e Tractian, projetei pipelines de IA local, RAG e orquestração de servidores MCP; "
                 "tive uma longa passagem de 6 anos pela IBM e XP como Arquiteto de Big Data; "
                 "e desenvolvi sistemas financeiros e de segurança de missão crítica para a B3 e o Itaú. "
                 "É essa bagagem de sistemas distribuídos e IA prática que trago consolidada para vocês hoje!"
    },

    # --- Abertura ---------------------------------------------------
    {
        "title": "AI sem Cloud com Golang!",
        "bullets": [
            "Foco em engenharia de sistemas de IA — Projetando e otimizando orquestração de modelos no próprio hardware",
            "Independência real — Como construir soluções robustas, seguras e de baixo custo fugindo de APIs proprietárias",
            "Infraestrutura local — O papel do desenvolvedor Go na consolidação de IA e dados locais com alta vazão"
        ],
        "notes": "Boa noite, pessoal! É um prazer enorme estar aqui no palco oficial da GolangSP. "
                 "A nossa talk de hoje aborda a IA sem Cloud. Vamos focalizar 100% em engenharia de sistemas de IA "
                 "e infraestrutura local-first. O nosso objetivo hoje é entender as engrenagens por dentro e provar "
                 "como Go é a linguagem ideal para construir esse motor local de alta resiliência e baixíssimo custo."
    },
    {
        "title": "O problema que ninguém fala em voz alta",
        "bullets": [
            "O ecossistema atual de agentes — cursor, cline, copilot: todos rodam em nuvem externa",
            "Custo oculto por clique — cada requisição consome tokens, enviando histórico acumulado",
            "Dependência de conectividade — se o Wi-Fi do evento ou do escritório falhar, o dev para",
            "Segurança de dados — código sensível trafegando para APIs de terceiros sem controle"
        ],
        "notes": "Quero começar com uma provocação. Quantos de vocês já "
                 "usaram um agente de IA pra programar? Quantos sabem "
                 "exatamente quanto cada interação custa, em tokens, em "
                 "dinheiro? É essa lacuna que essa talk ataca. Não é sobre "
                 "'IA é mágica'. É sobre entender o motor por dentro."
    },
    {
        "title": "CONTEXTO",
        "bullets": [
            "O custo do contexto — o tamanho do texto que você envia dita a velocidade e o preço",
            "O risco do ruído — contexto excessivo ou irrelevante confunde o raciocínio do modelo",
            "A grande pergunta de engenharia — como prover ao modelo exatamente o que ele precisa?",
            "A resposta moderna — MCP (Model Context Protocol) resolve a conexão padronizada"
        ],
        "notes": "Toda interação com um modelo de linguagem tem um custo "
                 "proporcional a uma coisa: quanto texto você manda pra "
                 "ele, o contexto. Quanto mais contexto irrelevante, mais "
                 "caro, mais lento, e às vezes pior a qualidade da "
                 "resposta."
    },
    {
        "title": "O que é o Protocolo MCP?",
        "bullets": [
            "Model Context Protocol — protocolo de comunicação aberto e padronizado criado pela Anthropic",
            "A analogia do USB da IA — você desenvolve um servidor MCP uma vez, qualquer agente o consome",
            "Padronização de ferramentas — define de forma unificada chamadas de ferramentas, prompts e recursos",
            "Go como porto seguro — escrever um servidor MCP é escrever um servidor HTTP/JSON estável e rápido"
        ],
        "notes": "MCP é um protocolo aberto que padroniza como um agente de "
                 "IA pede informação e executa ações fora de si mesmo. "
                 "Antes do MCP, cada ferramenta tinha sua forma proprietária "
                 "de se conectar a dados. O pulo do gato pra essa "
                 "comunidade Go: escrever um servidor MCP é escrever um "
                 "servidor, coisa que Go faz muito bem há muito tempo."
    },
    {
        "title": "Existe um jeito certo e um jeito ingênuo de fazer isso",
        "bullets": [
            "Jeito ingênuo (Cena 1) — carregar todo o repositório de dados de uma vez e enviar cru para a API",
            "Jeito inteligente (Cena 2) — realizar busca semântica local, enviando apenas o contexto ideal",
            "A infraestrutura necessária — rodar modelos de embeddings e inferência localmente via Kronk",
            "Conclusão empírica — comparativo real de custos e performance exibidos ao vivo no palco"
        ],
        "notes": "Vou mostrar dois jeito de construir um servidor MCP que "
                 "resolve o MESMO problema. Um é ingênuo: pega tudo, manda "
                 "tudo. O outro usa busca semântica de verdade, rodando "
                 "localmente, sem custar nuvem."
    },
    {
        "title": "Roteiro de hoje",
        "bullets": [
            "1. O motor local — mergulho técnico no Kronk (quantização, MoE, offload de CPU/GPU)",
            "2. Benchmarks reais — números de vazão, TTFT e latência medidos localmente",
            "3. Por dentro do código — walkthrough nos códigos Go dos dois servidores MCP",
            "4. Demonstração prática ao vivo — Cena 1 e Cena 2 duelando na tela com contadores ativos"
        ],
        "notes": "Esse é o mapa de hoje. Vamos primeiro entender o motor — "
                 "o Kronk. Depois abrir o código dos dois servidores MCP. "
                 "E só no final, com tudo entendido, ver os dois rodando "
                 "ao vivo, com o contador de tokens e custo do Cline na "
                 "tela. Bora?"
    },
    {
        "title": "O Motor de Toda IA Local: llama.cpp",
        "bullets": [
            "O divisor de águas da IA — Criado por Georgi Gerganov em 2023, reimplementando o LLaMA da Meta em C/C++ puro",
            "Independência absoluta de Python — Elimina a pesada pilha do PyTorch/Hugging Face por binários compilados nativos",
            "Aceleração de hardware multiplataforma — Suporte nativo a Metal (Apple Silicon), CUDA (Nvidia), ROCm (AMD) e CPU pura",
            "Pioneirismo no formato GGUF — Introduz arquivos quantizados ultra-eficientes mapeados diretamente em memória",
            "Mapeamento por FFI — Arquitetura ideal para ser importada nativamente por linguagens como o Go sem overhead"
        ],
        "notes": "Para entender como rodamos IA local de alta performance, precisamos falar do projeto open-source mais importante do ecossistema moderno: o llama.cpp. Criado por Georgi Gerganov no início de 2023, o seu objetivo fundamental era rodar o modelo LLaMA da Meta localmente usando C/C++ puro, sem depender da pesada pilha tradicional de Python (como PyTorch e Hugging Face). O llama.cpp compilado roda direto na máquina, falando diretamente com as otimizações vetoriais da sua CPU ou GPU. É por causa dele que a IA local-first é viável hoje de verdade. E como ele é escrito in C++, nós conseguimos nos conectar a ele diretamente a partir de linguagens como o Go!"
    },
    {
        "title": "O que é um Token de Verdade? Tokenização & BPE",
        "bullets": [
            "Fatiamento de Texto — Modelos de linguagem não leem palavras diretamente; eles quebram textos em fragmentos chamados tokens",
            "Byte Pair Encoding (BPE) — Algoritmo estatístico que agrupa pares de caracteres recorrentes para formar o vocabulário do modelo",
            "Mapeamento Numérico — Cada token é traduzido em uma ID numérica que aponta para um vetor de coordenadas na matriz do modelo",
            "Relação Espacial — 1 token equivale a aproximadamente 4 caracteres ou 0.75 palavras em inglês, variando de acordo com o idioma",
            "Tokens Especiais — Marcadores de estrutura essenciais: BOS (começo do prompt), EOS (fim da resposta) e PAD (preenchimento)",
            "Sensibilidade de Vocabulário — Termos comuns (ex: 'ChatGPT') viram 1 único token, mas palavras raras são divididas em 3 ou mais"
        ],
        "notes": "Para falarmos de custos e contadores de tokens na nossa demo, precisamos primeiro entender o que é um token de verdade na engenharia de LLMs. Modelos de inteligência artificial não sabem ler texto. Eles enxergam apenas números. O primeiro passo da inferência é a tokenização, onde um algoritmo estatístico chamado BPE (Byte Pair Encoding) quebra as palavras em fragmentos. A palavra 'ChatGPT', por exemplo, que é extremamente comum, é convertida em um único token pelo modelo. Mas uma palavra rara ou código Go complexo pode ser fatiado em 3 ou 4 tokens. Em média, no ecossistema, 1 token equivale a 4 caracteres, ou 0.75 palavras. Além disso, o modelo usa tokens especiais essenciais, como BOS para indicar o início da instrução, e EOS para indicar ao servidor que ele terminou de gerar e pode encerrar a conexão. É essa quebra estatística que define o volume de dados trafegado que o Cline mede ao vivo!"
    },
    {
        "title": "Além do Texto: O Ecossistema whisper.cpp",
        "bullets": [
            "Transcrição de Áudio Offline — Port de altíssima performance em C/C++ do modelo de reconhecimento de voz Whisper da OpenAI",
            "Mesma Filosofia de Engenharia — Desenvolvido também por Georgi Gerganov, eliminando dependências de Python por completo",
            "Velocidade de Processamento Extrema — Capaz de processar e transcrever áudio em até 10x mais rápido que o tempo real",
            "Suporte Amplo de Hardware — Otimizado nativamente via Metal, CUDA, CoreML e instruções de vetores de CPU",
            "Integração com Pipelines em Go — Comunicação direta via CGo assim como o Kronk, viabilizando agentes locais por voz"
        ],
        "notes": "E a genialidade desse ecossistema local-first não para no texto com o llama.cpp. O mesmo criador, Georgi Gerganov, desenvolveu outro projeto monumental: o whisper.cpp. Ele é uma reimplementação em C/C++ puro do modelo Whisper da OpenAI para reconhecimento de voz de altíssima precisão. Ele segue exatamente a mesma filosofia de engenharia: compila em um único executável nativo leve, roda totalmente offline, e é extremamente otimizado usando Metal no Mac e CUDA no PC. Ele consegue transcrever áudio em até 10 vezes mais rápido que o tempo real. E a conexão com Go é idêntica: você pode integrá-lo via CGo para construir pipelines de voz completos offline, onde seu agente Go escuta via whisper.cpp, raciocina com o Kronk, e pode até falar localmente. Isso prova que a IA offline-first não se limita a chat, mas sim a qualquer modalidade de mídia!"
    },
    {
        "title": "O Choque de Paradigmas: Python vs Go & C++",
        "table": [
            ["Aspecto", "Pilha Python (PyTorch / HF)", "Pilha Nativa (C++ / Go)"],
            ["Linguagem", "Python (Pesado, Interpretado, GIL)", "C/C++ e Go (Compilado, Concorrência Real)"],
            ["Runtime", "GigaBytes (PyTorch, bibliotecas, CUDA)", "MegaBytes (Executável compilado leve)"],
            ["Startup", "Lento (segundos p/ carregar interpretador)", "Instantâneo (milissegundos)"],
            ["RAM/VRAM", "Alto overhead de alocação de objetos", "Mapeamento direto em disco via mmap"],
            ["Foco Ideal", "Pesquisa acadêmica e treinamento", "Produção local, borda e servidores leves"]
        ],
        "notes": "Aqui está o choque de paradigmas entre a pilha tradicional e o que estamos propondo hoje. A pilha de pesquisa — PyTorch e Hugging Face — é incrível para criar e treinar modelos. Mas ela arrasta consigo o interpretador Python, gigabytes de dependências, startups lentos e um consumo de RAM absurdo de alocação de objetos. Quando migramos para llama.cpp compilado e o Kronk em Go, nós eliminamos o Python por completo da produção. O Go conversa diretamente com a biblioteca C++ do llama.cpp via CGo. Ganhamos a velocidade e o controle fino de hardware do C++ com a simplicidade, goroutines de alta concorrência e robustez de rede do Go!"
    },
    {
        "title": "Os Cérebros da Demo: Qwen & Gemma",
        "bullets": [
            "Modelos Open-Weights de elite — A liberação de modelos de ponta pelo Google e Alibaba viabilizou IA local corporativa",
            "Gemma (Google DeepMind) — Arquitetura de altíssima eficiência herdada diretamente da família comercial Gemini",
            "Mapeamento com embeddinggemma — Usamos o modelo Gemma de 300M parâmetros focado unicamente em embeddings semânticos",
            "Qwen (Alibaba Cloud) — Campeão absoluto de programação open-source, batendo GPT-4 em raciocínio de código",
            "Inferência com gpt-oss-20b — Nosso modelo principal de chat de 20B parâmetros, com MoE para acelerar em hardware local"
        ],
        "notes": "E quem são os cérebros que estão rodando de verdade nessa demo? São duas das melhores famílias de modelos de código aberto do mundo hoje: o Gemma do Google e o Qwen da Alibaba. Para a busca semântica, usamos o embeddinggemma-300m, que é um modelo ultra-leve do Google focado só em mapear textos em coordenadas matemáticas de 768 dimensões. E para resolver o bug em si, usamos o gpt-oss-20b, um modelo baseado no Qwen com 20 bilhões de parâmetros e Mixture of Experts que bate o GPT-4 da OpenAI em geração de código. É graças a esses gigantes liberando seus pesos de graça que podemos rodar IA de altíssimo nível local e de graça!"
    },
    {
        "title": "O Dilema de Seleção: Qwen/Gemma vs Alternativas",
        "table": [
            ["Modelo", "Arquitetura", "Vantagem Local-First", "Por que NÃO usamos para o Chat?"],
            ["Llama 3 (8B) / Mistral (7B)", "Denso (Dense)", "Populares no Mac e fáceis de carregar", "Menor precisão lógica para refatoração de concorrência sênior"],
            ["Microsoft Phi-3 (3.8B)", "Denso (Dense)", "Extremamente leve e velocidade altíssima", "Contexto curto e baixa inteligência estrutural para refatoração"],
            ["Google Gemma (Embed)", "Especializado", "Vetorização local ultra-veloz integrada no Kronk", "Modelo puramente de Embeddings (usado na Cena 2 para o Vault)"],
            ["Alibaba Qwen (MoE)", "Experts (MoE)", "Inteligência de 20B com custo de 4B por token", "Eleito para a Demo! Une precisão lógica e rapidez local-first"]
        ],
        "notes": "Aqui está a matriz de decisão de arquitetura para a escolha dos nossos modelos. É uma pergunta muito comum de equipes de tecnologia: 'Por que não usar o LLaMA 3 de 8B ou o Mistral que todo mundo conhece?'. A resposta é pura engenharia: modelos densos menores, como Llama 3 8B, são fáceis de carregar mas carecem da capacidade cerebral sênior necessária para identificar race conditions complexas e refatorar código concorrente em Go — eles alucinam com mais facilidade. O Phi-3 da Microsoft é super rápido, mas muito limitado para tarefas lógicas complexas. O Qwen-20B-MoE (nosso gpt-oss-20b) nos dá o melhor de dois mundos: a inteligência refinada de um modelo de 20 bilhões de parâmetros, mas com a latência e o custo computacional de um modelo de 4B de parâmetros porque ele é um Mixture of Experts (MoE) que ativa apenas 4 experts por token! E o Gemma de 300M é perfeito para embeddings semânticos, rodando na mesma memória compartilhada do Kronk em Go."
    },

    # --- Deep Dive: Kronk --------------------------------------------
    {
        "title": "Slide 1 — O que é o Kronk?",
        "bullets": [
            "Ardan Labs SDK — biblioteca e servidor nativo em Go projetado para envolver o motor llama.cpp",
            "Independência do Python — execute inferência e embeddings locais sem runtime pesado ou docker extra",
            "Compatibilidade OpenAI v1 — expõe endpoints padronizados permitindo troca de URL instantânea",
            "Robustez de infraestrutura — fornece pool de modelos ativos, cache incremental e monitoramento de logs"
        ],
        "notes": "Kronk é um SDK em Go, da Ardan Labs, que envolve o "
                 "llama.cpp. A grande sacada: ele expõe uma API compatível "
                 "com a da OpenAI. Kronk não é 'mais um wrapper de "
                 "chatbot'. É infraestrutura, pensada pra produção."
    },
    {
        "title": "Slide 2 — Arquitetura: onde o Kronk entra",
        "bullets": [
            "Camada do Agente (Cliente) — Cline, Cursor, Zed ou scripts Go conversando via HTTP OpenAI API",
            "Roteador Kronk Server — intercepta requisições, gerencia autenticação e gerencia o pool de modelos",
            "Motor Core llama.cpp — ponte CGo de alto desempenho para backends de hardware otimizados",
            "Aceleração de hardware nativa — seleciona automaticamente Metal (Mac), CUDA (Nvidia), ROCm ou CPU",
            "Formato unificado GGUF — modelos quantizados carregados de forma ultra-eficiente em disco"
        ],
        "notes": "O Kronk fica no meio. De um lado, qualquer cliente que "
                 "fale o protocolo da OpenAI. Do outro, ele escolhe "
                 "automaticamente o melhor backend de hardware disponível. "
                 "O modelo é um arquivo .gguf, formato do ecossistema "
                 "llama.cpp."
    },
    {
        "title": "Slide 3 — Por dentro do Kronk Server",
        "bullets": [
            "Roteador API principal — endpoint na porta :11435 para Chat Completions e Embeddings compatíveis",
            "Painel de inspeção de debug — monitoramento de métricas em tempo real e logs na porta :11445",
            "Mecanismo de Auth em produção — suporte integrado a JWT e chaves de API desativados para ambiente local",
            "Servidor MCP embutido — porta :9000 permitindo que agentes interajam diretamente com o orquestrador"
        ],
        "code_lang": "json",
        "code": '"status":"api router started","host":"0.0.0.0:11435"\n'
                '"status":"debug v1 router started","host":"0.0.0.0:11445"\n'
                '"status":"local mcp server started","host":"127.0.0.1:9000"\n'
                '"status":"start auth server"',
        "notes": "O Slide 2 mostrou o Kronk como uma caixa preta no meio do diagrama. Vamos abrir essa caixa rapidinho, com log real de hoje, não slide teórico. Quando você dá `kronk server start`, sobem pelo menos quatro coisas: o router da API (porta 11435, é o `/v1/chat/completions` compatível com OpenAI que o Cline fala), um router de debug (11445, métricas e introspecção), um serviço de auth (JWT/API keys — hoje desligado porque é uso local, mas existe pra produção), e — detalhe curioso — um servidor MCP nativo do próprio Kronk, na porta 9000. Não é nenhum dos nossos dois MCPs de demo; é o Kronk se expondo via MCP por conta própria, pra outros agentes consultarem o estado dele direto. O ponto de fundo: isso não é um wrapper fininho em volta do llama.cpp. É um processo de servidor com superfície de operação real — roteamento, auth, observabilidade — parecido com qualquer outro serviço backend que vocês já rodam em Go."
    },
    {
        "title": "Slide 4 — Pool de modelos",
        "bullets": [
            "Controle de orçamento de RAM/VRAM — limita por segurança o consumo para modelos (padrão: 80%)",
            "Proteção de OOM — recusa o carregamento se a estimativa estourar o orçamento, prevenindo quebras",
            "Multi-tenancy de modelos — suporta até 10 modelos residentes na memória simultaneamente",
            "Uso híbrido simultâneo — chat (ex: gpt-oss-20b) e embeddings (ex: embeddinggemma) ativos ao mesmo tempo",
            "TTL de desalocação — modelos ociosos são descarregados automaticamente para liberar memória do sistema"
        ],
        "code_lang": "json",
        "code": '"status":"resman-init","budget-percent":80,\n'
                '"ram-budget":"38.7GB","max-models-in-pool":10',
        "notes": "O Kronk não trata modelo como 'um arquivo que você abre uma vez'. Ele tem um pool de modelos — igual um pool de conexões de banco, só que pra modelos GGUF carregados na memória. Três coisas que esse pool faz por você: primeiro, orçamento de memória: por padrão, o Kronk só usa 80% da RAM/VRAM disponível pra modelos — aqui, hoje, isso virou um orçamento de 38.7GB. Se um modelo não cabe no orçamento, o Kronk recusa carregar antes de tentar em vez de só travar feio na hora de alocar — foi exatamente esse mecanismo de segurança que pegou a configuração original dessa talk, preparada num Linux com GPU dedicada de 12GB, antes da gente migrar pra esse Mac com memória unificada (mais sobre isso no Slide 7, quando chegarmos em MoE). Segundo, múltiplos modelos ao mesmo tempo: o pool aceita até 10 modelos residentes simultaneamente. Não é teórico — agora mesmo, nesta máquina, tem dois modelos carregados ao mesmo tempo: o gpt-oss-20b pra chat, e o embeddinggemma pra busca semântica do vault_smart. Cada um com seu próprio orçamento de memória, dividindo o mesmo processo Kronk. Terceiro, TTL de ociosidade: modelo sem uso por um tempo configurável é descarregado automaticamente, liberando memória pra outra coisa — sem você precisar gerenciar isso manualmente."
    },
    {
        "title": "Slide 5 — Cache incremental (IMC)",
        "bullets": [
            "O padrão de agentes de IA — conversas longas reenviando todo o histórico acumulado a cada turno",
            "Gargalo sem cache — reprocessar o prompt inteiro consome ciclos e atrasa o tempo de resposta",
            "Cache Incremental de Mensagens (IMC) — retém o KV-cache internamente na memória entre chamadas",
            "Aceleração instantânea — processa apenas os tokens de fato novos, reduzindo drasticamente o TTFT"
        ],
        "code_lang": "yaml",
        "code": "incremental-cache: false   # Incremental message caching\n"
                "                           # for agentic workflows\n"
                "                           # (Cline, OpenCode)\n"
                "# valor real hoje, confirmado no log: IncrementalCache[true]",
        "notes": "Uma conversa com o Cline reenvia todo o histórico "
                 "anterior a cada turno. O cache incremental de mensagens "
                 "guarda o estado da conversa e só processa os tokens "
                 "novos no turno seguinte. Vem ligado por padrão pra "
                 "clientes agênticos."
    },
    {
        "title": "Por trás da Inferência: O que é o KV-Cache?",
        "bullets": [
            "O Custo de Auto-Atenção — Em redes baseadas em Transformers (LLMs), a geração de cada novo token exige analisar todo o histórico anterior",
            "Complexidade Quadrática — Sem otimização, computar matrizes de atenção para prompts longos consome processamento absurdo de hardware",
            "Keys & Values (KV) — O modelo salva as matrizes de Chaves (K) e Valores (V) dos tokens processados direto na VRAM",
            "Geração Linear — Evita recalcular o passado; novos tokens são computados aproveitando o KV-Cache de forma instantânea",
            "Impacto com Agentes — Conversas de ~6k tokens caem de 12 segundos de espera para menos de 0.5 segundos com o IMC do Kronk"
        ],
        "notes": "Para uma audiência sênior de engenharia como a nossa hoje, explicar o KV-Cache é abrir as entranhas dos Transformers. LLMs geram texto token por token de forma autorregressiva. Para gerar a próxima palavra, ele precisa olhar para tudo o que veio antes. Na auto-atenção, cada token tem vetores de Query (pergunta), Key (endereço) e Value (conteúdo). Se não tivéssemos o KV-Cache, o modelo teria que recalcular as Chaves e Valores de todo o histórico a cada novo token gerado — uma complexidade quadrática absurda! O KV-Cache salva esses vetores K e V diretamente na VRAM. O modelo só calcula o QKV para o único token novo gerado e consulta o passado direto no cache. É esse buffer físico que o Cache Incremental (IMC) do Kronk mantém quente na memória entre os turnos do Cline, fazendo o tempo até o primeiro token de prompts gigantes despencar de 12 segundos para menos de meio segundo local-first!"
    },
    {
        "title": "Slide 6 — Quantização (a base de tudo)",
        "table": [
            ["Formato", "Bits", "Consumo RAM", "Perda de Precisão"],
            ["F16", "16", "Muito Alto (~24GB)", "Nenhuma (Precisão nativa)"],
            ["Q8_0", "8", "Médio (~12GB)", "Quase Nula (Excelente p/ código)"],
            ["Q4_K_M", "4", "Baixo (~6.5GB)", "Pequena (Leve perda semântica)"]
        ],
        "bullets": [
            "Nosso setup hoje — modelo gpt-oss-20b quantizado em Q8_0, ocupando ~12GB estáveis em disco",
            "O ganho prático — redução de mais de 50% de RAM sem impactos na acurácia lógica dos testes"
        ],
        "notes": "Quantização é comprimir os números de ponto flutuante "
                 "do modelo pra ocupar menos memória. Q8_0 usa 8 bits — "
                 "praticamente sem perda. Q4_K_M usa 4 bits — mais "
                 "compacto, perdinha quase imperceptível em código."
    },
    {
        "title": "Slide 7 — MoE (Mixture of Experts)",
        "bullets": [
            "Arquitetura modular MoE — modelo fatiado em sub-redes especialistas em vez de um bloco monolítico",
            "Ativação dinâmica — gpt-oss-20b possui 20B parâmetros e 32 especialistas, mas ativa apenas 4 por token",
            "Suporte nativo do Kronk — moe.mode permite alocar roteador e especialistas de forma inteligente",
            "Políticas de distribuição — auto, experts_cpu (expert na CPU, roteador na GPU), experts_gpu e keep_top_n",
            "Relação de hardware — no Linux com GPU de 12GB estourávamos a VRAM; no Mac M4 Pro roda suave unificado"
        ],
        "notes": "MoE: para cada token, um roteador decide quais "
                 "especialistas ativar — só uma fração do total. O "
                 "gpt-oss-20b tem 20B parâmetros mas custo computacional "
                 "de um modelo bem menor. Nota de bastidor: no Linux com "
                 "GPU de 12GB, o gpt-oss-20b não cabia com contexto "
                 "grande — usamos um modelo denso menor. No Mac, com "
                 "memória unificada, o MoE de 20B cabe com folga."
    },
    {
        "title": "Por dentro da Arquitetura MoE (Mixture of Experts)",
        "bullets": [
            "Fatiamento de Rede (FFNs) — Camadas neurais tradicionais divididas em sub-redes especialistas independentes",
            "Portão de Roteamento (Gating Network) — Classificador ultra leve calcula e ordena afinidades de cada token",
            "Inferência Top-K — Apenas os N melhores especialistas (ex: 4 de 32) são acionados na computação por token",
            "Computação vs Memória — 20B parâmetros residentes na RAM, mas com custo computacional (FLOPs) de apenas 4B",
            "Vantagem Local-First — Kronk moe.mode distribui especialistas entre CPU/GPU se VRAM for gargalo"
        ],
        "notes": "Para entender Mixture of Experts, pensem em um hospital: em vez de consultar todos os médicos da instituição para cada dor de cabeça (que seria um modelo denso tradicional acendendo todos os seus bilhões de neurônios), um roteador na recepção direciona você diretamente para os dois ou três especialistas ideais. No gpt-oss-20b que usamos hoje, o modelo possui 20 bilhões de parâmetros na memória, mas para cada token processado, ele ativa apenas 4 especialistas. Isso garante a acurácia de um cérebro gigante com a velocidade e o consumo computacional de um modelo pequeno de 4B. E o Kronk nos dá políticas para colocar esses especialistas na CPU ou GPU de forma transparente, permitindo que MoEs massivos rodem localmente em hardware comum."
    },
    {
        "title": "Slide 8 — Hybrid CPU/GPU (offload de camadas)",
        "bullets": [
            "Configuração ngpu-layers — controla rigorosamente a quantidade de camadas de transformer na GPU",
            "Pegadinha de convenção do Kronk — valor 0 significa TODAS na GPU, -1 força CPU pura, N move N camadas",
            "Gerenciamento fino de VRAM — dial essencial em GPUs dedicadas para balancear o cache de contexto",
            "O cenário da unificada — no Apple Silicon a RAM é compartilhada, operamos ngpu-layers:0 sem gargalos"
        ],
        "notes": "Pegadinha real que vivemos: a convenção do Kronk é "
                 "invertida do que se imagina — 0 significa todas as "
                 "camadas na GPU. Numa GPU dedicada isso é um dial real "
                 "entre velocidade e memória; no Mac com memória "
                 "unificada, o gargalo simplesmente não existe da mesma "
                 "forma."
    },
    {
        "title": "Slide 9 — Nossos números reais (hoje, nesta máquina)",
        "table": [
            ["Configuração", "Prompt tokens", "TTFT (Time-to-first-token)", "Tokens/s"],
            ["CPU only", "pequeno (~100)", "—", "~6-8 tok/s"],
            ["Metal full offload — curto", "81", "~0.2 - 0.4s", "~70 tok/s"],
            ["Metal full offload — médio (Cline real)", "2.034", "~2.2s", "~65-68 tok/s"],
            ["Metal full offload — stress (MCPs grandes)", "9.084", "~11.2 - 12.4s", "~63 tok/s"]
        ],
        "bullets": [
            "Vantagem do hardware local — aceleração do Metal chega a ser de 9x a 10x mais rápida que CPU pura",
            "Estabilidade sob estresse — contextos massivos de ~9k tokens não provocam estouros ou falhas"
        ],
        "notes": "MacBook M4 Pro, 48GB de memória unificada. Metal full "
                 "offload é uns 9-10x mais rápido que CPU pura. E nenhuma "
                 "linha falha por falta de memória — vantagem real da "
                 "memória unificada. Mas teve uma pegadinha, e ela ensina "
                 "uma lição melhor que qualquer slide teórico."
    },
    {
        "title": "Slide 10 — Estudo de Caso: Análise de Gargalos de Prefill",
        "bullets": [
            "Identificação da latência — observamos TTFT de ~12s no prompt longo de estresse e assumimos gargalo de IO",
            "Hipótese inicial — presumimos que o tamanho de batch padrão (nbatch/nubatch) estava estrangulando o Metal",
            "Tentativa de otimização — aumentamos e quadruplicamos os buffers na configuração do YAML do Kronk",
            "Resultado surpreendente — ganho irrisório de 6% (latência caindo minimamente de 12.36s para 11.58s)",
            "Lição de engenharia — o gargalo era computacional (roteamento de especialistas MoE), não transferência de dados. Meça!"
        ],
        "notes": "Vendo o TTFT de ~12s, nosso instinto foi 'aumentar o "
                 "batch'. Resultado: 6% de ganho, quase nada. A lição "
                 "real: pra esse modelo MoE no Metal, o prefill é "
                 "limitado por compute (roteamento de especialistas), não "
                 "por throughput de batch. Só descobriu porque mediu "
                 "antes de assumir."
    },
    {
        "title": 'Slide 11 — "Thinking mode" e o custo do raciocínio',
        "bullets": [
            "Modo de raciocínio lógico — gpt-oss-20b emite blocos extensos de cadeia de pensamento antes da resposta",
            "Exemplo do oi — uma simples saudação de teste consumiu 40 tokens de pensamento oculto antes de responder",
            "Parâmetro de requisição — controlado em tempo de execução via reasoning_effort (low / medium / high)",
            "Tradeoff clássico de IA — esforço alto garante correção precisa de algoritmos, mas penaliza a latência"
        ],
        "notes": "O gpt-oss-20b tem um modo de raciocínio que melhora a "
                 "qualidade da decisão, mas cada token de pensamento "
                 "custa tempo. O Kronk deixa controlar isso por "
                 "requisição, com reasoning_effort — não é configuração "
                 "fixa do modelo."
    },
    {
        "title": "Slide 12 — Trocar de modelo é runtime, não startup",
        "code_lang": "bash",
        "code": "kronk model list --local\nkronk model pull <nome>\n\n# depois, só aponta o cliente pro novo nome:\n# Cline -> Settings -> Model ID\n# curl  -> campo \"model\" no JSON",
        "bullets": [
            "O funcionamento do pool — o boot do servidor apenas prepara o orquestrador, mantendo o pool vazio",
            "Escolha dinâmica por chamada — o modelo é lido sob demanda com base no campo 'model' no JSON",
            "Carregamento lazy-load — se o modelo não estiver quente, o Kronk o monta e descarrega quando ocioso",
            "Sem restarts de infra — altere de Qwen para Gemma em tempo real de execução mudando apenas as chamadas"
        ],
        "notes": "Detalhe que confunde quem vem de outras stacks: você "
                 "não escolhe o modelo no boot do servidor. O pool "
                 "decide por requisição, igual a API da OpenAI. Pra "
                 "trocar: confere com model list --local, baixa com "
                 "model pull se precisar, e só aponta o cliente pro "
                 "novo nome — sem reiniciar o kronk server start."
    },
    {
        "title": "Slide 13 — Maior ou menor? O dial de tamanho no model_config.yaml",
        "table": [
            ["Modelo", "Tamanho / Tipo", "Contexto Suportado", "Status Atual"],
            ["Qwopus3.5-4B-Coder", "4B (Ultra leve)", "73.728 tokens", "Não Baixado (configurado)"],
            ["gpt-oss-20b-Q8_0", "20B (MoE, Equilibrado)", "32.768 tokens", "Baixado e Quente (Uso atual)"],
            ["gemma-4-26B-A4B-it", "26B (MoE, Completo)", "131.072 tokens", "Não Baixado (configurado)"],
            ["Qwen3.6-35B-A3B", "35B (MoE, Avançado)", "131.072 tokens", "Não Baixado (configurado)"]
        ],
        "bullets": [
            "Flexibilidade por hardware — dial de modelos prontos pré-configurados no arquivo local do usuário",
            "Ajuste fino de escopo — modelos de 4B operam rápido em notebooks modestos, MoEs maiores escalam robustez"
        ],
        "notes": "Tamanho de modelo também é um dial, não decisão "
                 "travada em pedra. Mesma API, só muda o nome no campo "
                 "model (e baixar antes, se for novo). Já tenho tuning "
                 "pronto pra dois modelos bem maiores que o de hoje e "
                 "um menor (4B) pra hardware mais limitado — a config "
                 "já existe no model_config.yaml desta máquina, só "
                 "falta o pull."
    },
    {
        "title": "Slide 14 — Resumo / Takeaways",
        "bullets": [
            "Local-first viável — Kronk com aceleração nativa Metal entrega vazão comercializável em desktop local",
            "Conceito de MoE consistente — especialistas segmentados economizam processamento em qualquer hardware",
            "Documentação manda — a escala invertida de ngpu-layers do Kronk exige atenção para não travar em CPU",
            "Empirismo na prática — otimização de IA exige medição de bottlenecks reais, não chutes",
            "Flexibilidade de runtime — controle de raciocínio e troca de modelos ao vivo na infra em execução"
        ],
        "notes": "É exatamente essa flexibilidade — CPU, GPU, híbrido, "
                 "MoE, raciocínio ligado ou desligado — que torna o "
                 "Kronk uma peça de infraestrutura séria pra rodar IA "
                 "local em produção, não só em demo de palco."
    },
    {
        "title": "A Grande Pergunta: Kronk vs Ollama",
        "bullets": [
            "Por que não o Ollama? — Todo desenvolvedor Go usará essa pergunta clássica como barreira de entrada técnica",
            "Ollama (A Abordagem Daemon) — Excelente ferramenta CLI desktop, mas atua como daemon fechado externo e opaco",
            "Kronk (In-Process Go SDK) — Biblioteca nativa em Go que permite carregar modelos e embeddings diretamente no seu processo Go",
            "Embutido de Verdade — Usamos o Kronk SDK dentro do `vault_smart` para carregar o `embeddinggemma` de forma programática",
            "Otimizações de Agentes (IMC) — O Kronk traz o Cache Incremental de Mensagens (IMC) nativo e afinado para workflows agênticos",
            "Controle de Baixo Nível — Controle total sobre buffers, threads, caches e pool de modelos declarados via YAML ou código Go"
        ],
        "notes": "Essa é a grande pergunta que qualquer desenvolvedor de Go na plateia vai fazer: 'Por que vocês estão usando o Kronk se já temos o Ollama?'. A resposta é a diferença entre um executável de prateleira e uma biblioteca de engenharia. O Ollama é fantástico para uso desktop local, mas ele roda como um daemon fechado externo. Você não consegue embutir o Ollama dentro do seu binário compilado em Go. Já o Kronk é um SDK Go nativo FFI CGo. Na nossa Cena 2, o servidor `vault_smart` utiliza o Kronk SDK para carregar e instanciar programaticamente o modelo de embeddings `embeddinggemma` dentro do mesmo processo em execução em Go! Além disso, o Kronk traz otimizações agênticas nativas pensadas para o Cline (como o IMC) e controle total de baixo nível de memória, threads e pool de modelos que o Ollama oculta de forma opaca."
    },
    {
        "title": "O Desafio do Contexto: Por que 25 Tasks?",
        "bullets": [
            "Ruído Estrutural do Backlog — Projetos reais contêm centenas de tickets, arquivos legados e documentações",
            "Inchaço de Janela de Contexto — Abordagem ingênua (manda tudo) consome ~47k tokens por chamada",
            "Diluição de Atenção (Lost in the Middle) — Modelos tendem a ignorar detalhes ocultos no meio de prompts gigantes",
            "Degradação de Velocidade — Processar 47k de histórico local-first eleva o tempo de prefill inviabilizando a resposta",
            "Inviabilidade Financeira — Enviar ruído para APIs pagas destrói a viabilidade comercial do projeto"
        ],
        "notes": "Na nossa demonstração, o Obsidian possui 25 tarefas antigas e apenas uma tarefa relevante com o bug real, a TASK-001. Por que fizemos isso? Porque em um cenário corporativo real, o repositório do time tem centenas de páginas de documentação. Se usarmos a abordagem ingênua da Cena 1 de carregar tudo e mandar pro modelo de uma vez, estamos enviando 47 mil tokens de puro ruído. O modelo sofre de diluição de atenção — um fenômeno conhecido na literatura onde LLMs esquecem detalhes no meio do prompt — e a velocidade despenca de forma inviável. Precisamos filtrar antes de enviar."
    },
    {
        "title": "A Solução Inteligente: Busca Semântica",
        "bullets": [
            "Embeddings de Documentos — Modelo embeddinggemma converte cada arquivo markdown em um vetor numérico denso",
            "Coordenadas de Significado — O vetor de 768 dimensões mapeia a semântica do texto em um espaço geométrico",
            "Busca por Proximidade — A pergunta do desenvolvedor é vetorizada e comparada geometricamente com o backlog",
            "Superação de Palavras-Chave — Termos correlatos (ex: 'faturamento' e 'reembolso') aproximam-se sem termos idênticos",
            "Filtragem Cirúrgica — Recupera apenas as top-3 tarefas mais compatíveis, derrubando o contexto em 90%"
        ],
        "notes": "Como resolvemos o problema do ruído? Usamos busca semântica de verdade. Antes do agente fazer qualquer pergunta, um modelo de embedding local lê os 26 markdown do Obsidian e transforma cada um em um vetor matemático de 768 coordenadas. Quando o agente Cline faz a pergunta, nós convertemos essa pergunta no mesmo espaço de coordenadas. Através da matemática de vetores, calculamos quais documentos estão geometricamente mais próximos da pergunta. O resultado é cirúrgico: o Smart MCP envia apenas os 3 documentos mais relevantes. Saímos de 47 mil tokens para meros 4 mil tokens de contexto, eliminando o ruído e mantendo a velocidade máxima de resposta."
    },
    {
        "title": "O Motor do Vetor: DuckDB + VSS + HNSW",
        "bullets": [
            "Mecanismo In-Memory — DuckDB inicializa instantaneamente em RAM dentro do próprio processo Go do MCP",
            "Extensão DuckDB VSS — Adiciona suporte nativo a tipos vetoriais (FLOAT[768]) e funções matemáticas",
            "Índice HNSW (Hierarchical Navigable Small World) — Algoritmo de grafo para busca rápida de vizinhos próximos",
            "Busca em Tempo Logarítmico — Evita varredura linear lenta, escalando perfeitamente para milhões de registros",
            "Similaridade de Cosseno via SQL — Query simples em banco relacional calcula e ordena os vetores por relevância"
        ],
        "notes": "E como implementamos essa busca em Go sem complexidade de infraestrutura de banco de dados corporativa externa? Usamos o DuckDB rodando em memória direto no processo Go do MCP! Com a extensão de Busca de Similaridade Vetorial (VSS), conseguimos criar um índice HNSW — que é uma estrutura de grafos extremamente rápida para vizinhos mais próximos. A busca vira uma query SQL relacional simples que calcula a similaridade de cosseno entre o vetor da pergunta e os vetores indexados das tarefas do Obsidian. Isso acontece em milissegundos e nos dá o top-3 de forma extremamente eficiente."
    },
    {
        "title": "Chunk vs Embedding: O Pipeline do RAG",
        "bullets": [
            "O que é um Chunk? — Partição física de texto bruto legível (humano). Cada um dos 26 markdowns do Obsidian agirá como um chunk",
            "O que é um Embedding? — Tradução do significado semântico do chunk em uma coordenada matemática (vetor) de 768 floats",
            "A Conexão no RAG — DuckDB usa o embedding para calcular a proximidade e recuperar o chunk associado em milissegundos",
            "Economia de Contexto — O modelo de chat recebe o chunk de texto puro ideal, eliminando 90% do ruído do backlog corporativo"
        ],
        "code_lang": "text",
        "code": "[FASE 1: INDEXAÇÃO (No Startup)]\n"
                "Documento.md ───> CHUNK (Texto Puro)\n"
                "                        │\n"
                "                 (embeddinggemma)\n"
                "                        ▼\n"
                "                EMBEDDING (768 floats)\n"
                "                        │\n"
                "                        ▼\n"
                "              [ Salva no DuckDB ]\n"
                "\n"
                "[FASE 2: RECUPERAÇÃO (Na Pergunta)]\n"
                "Pergunta ──> Vetor Query ──> SQL Cosseno\n"
                "                                 │\n"
                "                                 ▼\n"
                "                           Acha EMBEDDING\n"
                "                                 │\n"
                "                                 ▼\n"
                "                           Retorna CHUNK",
        "notes": "Esta é a conceituação mais importante do RAG local: a diferença entre Chunk e Embedding. O Chunk é a fatia física do texto bruto — no nosso caso, o texto de cada uma das 26 tasks do Obsidian. O Embedding é a representação matemática do significado desse texto, gerado pelo modelo embeddinggemma como um vetor de 768 dimensões. O banco DuckDB usa o Embedding para fazer cálculos geométricos ultra-rápidos de cosseno e achar o endereço físico do texto. O resultado disso é que entregamos ao modelo de chat apenas o texto puro (Chunk) do ticket relevante, gerando economia máxima de tokens e alta performance local-first."
    },

    # --- MCP Walkthrough --------------------------------------------
    {
        "title": "Slide 15 — A base compartilhada: lendo o vault",
        "code_lang": "go",
        "code": 'func LoadAll(root string) ([]Document, error) {\n'
                '    var docs []Document\n'
                '    filepath.WalkDir(root, func(path string, d os.DirEntry,'
                ' err error) error {\n'
                '        if d.IsDir() || !strings.HasSuffix(d.Name(), '
                '".md") {\n'
                '            return nil\n'
                '        }\n'
                '        content, _ := os.ReadFile(path)\n'
                '        rel, _ := filepath.Rel(root, path)\n'
                '        docs = append(docs, Document{Path: rel, '
                'Content: string(content)})\n'
                '        return nil\n'
                '    })\n'
                '    return docs, nil\n'
                '}',
        "bullets": [
            "Lógica unificada — os dois servidores MCP de demo consomem exatamente a mesma função interna",
            "Escaneamento simples — percorre a pasta recursivamente coletando o caminho relativo e texto bruto",
            "Escala do teste — processa os 26 arquivos em Markdown simulando o backlog fictício",
            "O divisor de águas — o que dita a eficiência não é o carregamento, mas o filtro de busca"
        ],
        "notes": "LoadAll anda pela pasta do vault, pega todo .md, lê o "
                 "conteúdo bruto e devolve uma lista de documentos. "
                 "Simples, sem mágica. A diferença entre as duas cenas "
                 "não está em COMO os dados são lidos."
    },
    {
        "title": 'Slide 16 — Cena 1: Servidor MCP sem Filtragem (vault_dumb)',
        "code_lang": "go",
        "code": 'mcp.AddTool(server, &mcp.Tool{\n'
                '    Name: "read_vault",\n'
                '    Description: "Returns the full contents of every '
                'Markdown file...",\n'
                '}, readVaultHandler(vaultPath))\n\n'
                'func readVaultHandler(vaultPath string) ... {\n'
                '    docs, _ := vault.LoadAll(vaultPath)\n'
                '    // concatena tudo: --- FILE: path ---\\ncontent\\n\\n\n'
                '}',
        "bullets": [
            "Uma única ferramenta unificada — read_vault exposta sem necessidade de argumentos de entrada",
            "Processamento sem filtro — carrega todas as tarefas, concatena em um payload cru e devolve de uma vez",
            "Ignorância de relevância — o agente pede contexto de uma única task, e recebe o backlog inteiro",
            "Custo de simplicidade — apenas 20 linhas de desenvolvimento Go, resultando em ~47k tokens na nuvem"
        ],
        "notes": "O agente pediu 'me ajuda com a TASK-001' e essa "
                 "ferramenta devolve as 26 tasks completas. É proposital: "
                 "representa o jeito ingênuo de integrar LLM com dados. "
                 "É exatamente essa simplicidade que custa caro."
    },
    {
        "title": "Slide 17 — Cena 2, parte 1: montando o índice vetorial",
        "code_lang": "go",
        "code": 'krnEmbed, _ := kronk.New(model.WithModelFiles(\n'
                '    []string{embedModelPath}))\n'
                'docs, _ := vault.LoadAll(vaultPath)\n'
                'dims, _ := embeddingDims(krnEmbed)\n'
                'db, _ := store.LoadVault(ctx, krnEmbed, docs, dims)',
        "bullets": [
            "Modelo secundário local — carrega o embeddinggemma-300m, especialista veloz em mapear significados",
            "Representação matemática — gera vetores de 768 floats (coordenadas) representando semântica dos arquivos",
            "Banco de dados embutido — inicializa DuckDB localmente acoplado com índice HNSW para vizinhos próximos",
            "Mapeamento único — a vetorização e inserção ocorrem somente na inicialização do servidor, sem delay na requisição"
        ],
        "notes": "Antes de aceitar qualquer pergunta, a Cena 2 carrega um "
                 "modelo de embedding e transforma cada documento num "
                 "vetor — 'coordenadas de significado' num espaço "
                 "matemático. Isso vai pro DuckDB com índice HNSW."
    },
    {
        "title": "Slide 18 — Cena 2, parte 2: a busca de verdade",
        "code_lang": "sql",
        "code": "SELECT path, text,\n"
                "       array_cosine_similarity(embedding, "
                "?::FLOAT[768]) AS similarity\n"
                "FROM items\n"
                "ORDER BY similarity DESC\n"
                "LIMIT 3;",
        "bullets": [
            "Vetorização de query — transforma a pergunta do usuário no mesmo formato de coordenadas numéricas",
            "Busca vetorial via SQL — executa query no DuckDB calculando similaridade de cosseno com os vetores salvos",
            "Busca por significado — encontra termos análogos (ex: faturamento recupera cupom) sem palavras idênticas",
            "Retorno ultra enxuto — reduz o volume de texto enviado do repositório inteiro para apenas os top-3 úteis"
        ],
        "notes": "search_vault transforma a pergunta no mesmo tipo de "
                 "vetor e busca os 3 documentos mais similares por "
                 "cosseno. É por isso que perguntar sobre 'dinheiro e "
                 "pagamentos' acha a task de cupom mesmo sem a palavra "
                 "'dinheiro' aparecer."
    },
    {
        "title": "Slide 19 — Lado a lado: mesma pergunta, dois caminhos",
        "table": [
            ["Métrica / Recurso", "vault_dumb (Cena 1)", "vault_smart (Cena 2)"],
            ["Nome da Ferramenta", "read_vault", "search_vault(query, top_k)"],
            ["Mecanismo Interno", "Concatenação crua", "Embeddings + Busca Vetorial"],
            ["Armazenamento", "Memória volátil simples", "DuckDB em RAM + Índice HNSW"],
            ["Preparação no Boot", "Nenhuma", "Mapeamento semântico no startup"],
            ["Tokens Consumidos", "~47.000 tokens (alto custo)", "~4.000 tokens (alta otimização)"],
            ["Complexidade de Código", "Mínima (~20 linhas Go)", "Moderada (~80 linhas Go)"]
        ],
        "notes": "Cena 1 é 4x mais simples de código, mas custa 10x mais "
                 "tokens. Cena 2 exige infraestrutura extra, mas paga "
                 "isso de volta multiplicado em economia de contexto — e "
                 "tudo local, via Kronk, sem custar um centavo."
    },
    {
        "title": "O Alvo do Fix: Concorrência no DeductStock",
        "bullets": [
            "Código de E-Commerce — O pacote de estoque `internal/inventory/service.go` possui a função DeductStock",
            "Vulnerabilidade Concorrente — Acesso concorrente ao mapa Go `products` sem travas causa pânico fatal no runtime",
            "A Race Condition de Estado — Múltiplas goroutines podem validar o estoque simultaneamente antes que a subtração aconteça",
            "O Teste de Concorrência (`-race`) — Temos um teste de WaitGroup simulando 20 requisições simultâneas para um estoque de 10",
            "O Objetivo da Demo — O agente Cline precisa ler a TASK-001 via MCP, entender o bug de concorrência e aplicar sync.Mutex"
        ],
        "code_lang": "go",
        "code": "func (s *Service) DeductStock(productID string, quantity int) error {\n"
                "    p, ok := s.products[productID] // Concorrência no Mapa (Fatal!)\n"
                "    if !ok {\n"
                "        return fmt.Errorf(\"produto %s não encontrado\", productID)\n"
                "    }\n"
                "\n"
                "    // BUG: Race condition aqui. Múltiplas goroutines compram juntas!\n"
                "    if p.Stock < quantity {\n"
                "        return ErrInsufficientStock\n"
                "    }\n"
                "\n"
                "    p.Stock -= quantity // Modificação Concorrente sem Locks!\n"
                "    return nil\n"
                "}",
        "notes": "Antes de darmos o play nas demos ao vivo, deixa eu mostrar qual é o alvo físico que o nosso agente de IA precisa consertar na nossa base de código de teste. Temos um serviço simples de estoque em Go. Na função DeductStock, temos um bug de concorrência clássico e gravíssimo: primeiro, acessamos concorrentemente o mapa s.products sem nenhum lock (o que faz o Go explodir com fatal panic no runtime); segundo, decrementamos o estoque sem sincronização, permitindo que requisições simultâneas comprem mais do que temos em estoque (um clássico data race). Temos um teste concorrente preparado sob estresse com sync.WaitGroup que falha consistentemente ao rodar go test -race. O desafio do agente local Cline é consultar o Obsidian para descobrir o que a TASK-001 pede, encontrar esse arquivo, entender o bug de concorrência, aplicar uma trava de Mutex, e validar rodando go test com a flag -race habilitada. Vamos ver como a Cena 1 e a Cena 2 lidam com isso!"
    },
    {
        "title": "Como o Agente Pensa: O Loop ReAct",
        "bullets": [
            "Além de Chat Completions — Agentes não apenas geram texto; eles executam ações de forma iterativa",
            "Padrão ReAct (Reasoning & Acting) — Paradigma de execução que combina raciocínio lógico e chamadas de ferramentas",
            "Fase 1: Thought (Raciocínio) — O agente analisa o estado atual da tarefa e decide o que fazer",
            "Fase 2: Action (Ação) — Chama uma ferramenta do MCP (como `read_vault` ou `search_vault`)",
            "Fase 3: Observation (Observação) — O agente lê o resultado retornado pelo MCP e avalia o que mudou",
            "Ciclo de Resolução — O loop repete-se de forma autônoma até o bug ser encontrado, corrigido e validado com testes"
        ],
        "notes": "Antes de abrirmos a tela para as demos, precisamos entender uma coisa fundamental sobre o Cline: ele não é um chat simples onde você joga uma pergunta e ele cospe uma resposta. O Cline é um agente autônomo baseado no padrão de arquitetura ReAct — Reasoning and Acting. Esse padrão roda em um ciclo iterativo de três fases: primeiro, Thought — o agente escreve um monólogo interno pensando sobre o que precisa fazer; segundo, Action — ele decide chamar uma ferramenta externa, que no nosso caso é uma ferramenta MCP do nosso servidor; terceiro, Observation — ele lê o resultado retornado pelo MCP e avalia se conseguiu o que queria. Ele vai rodar esse loop de forma autônoma: vai ler o backlog Obsidian, achar o arquivo Go, ler o código, aplicar o Mutex, abrir o terminal, rodar go test -race e parar apenas quando os testes estiverem passando. É um processo de engenharia ativo executado por software!"
    },
    {
        "title": "O Protagonista da Demo: Configurando o Cline",
        "bullets": [
            "Extensão VS Code — O Cline opera como uma extensão de chat e automação de arquivos integrada no editor",
            "Apontamento Local — Configuramos o provedor para API Compatível com OpenAI, base URL `http://localhost:11435/v1`",
            "Seleção de Modelo — Apontamos o ID do modelo para `gpt-oss-20b-Q8_0` que o Kronk mantém carregado no pool",
            "Contador de Tokens — Cline exibe nativamente a quantidade de tokens consumida por turno e o custo em dólares",
            "O Clímax das Cenas — Na Cena 1, o MCP Dumb inunda o contexto (Custo: $1.43); na Cena 2, o MCP Smart filtra cirurgicamente (Custo: $0.00)"
        ],
        "notes": "O protagonista físico da nossa demo é o Cline, que roda aqui como uma extensão integrada ao VS Code. Ele tem acesso de leitura e escrita de arquivos e execução de terminal local. Para a Cena 1, nós o configuramos apontando para a API de nuvem tradicional com o MCP Dumb ativo. Para a Cena 2, nós mudamos o provider para apontar para o nosso Kronk local na porta 11435, com o modelo gpt-oss-20b e o MCP Smart ativo. O Cline tem um recurso sensacional que é o contador de tokens e custos embutido na barra de título do chat. A cada ferramenta chamada ou turno, ele atualiza na hora o gasto acumulado. É nessa caixinha cinza do Cline que o clímax da nossa talk acontece: vocês vão assistir ao vivo o custo de contexto na nuvem da Cena 1 bater mais de 1 dólar em uma única rodada, e o custo local da Cena 2 cravar zero absoluto!"
    },

    # --- Demos Finais -------------------------------------------------
    {
        "title": "Slide 20 — Chegou a hora",
        "bullets": [
            "Cenários idênticos — mesmo agente Cline, mesma pergunta disparada pelo prompt",
            "Abordagens antagônicas — busca semântica inteligente contra dump de dados massivos",
            "Dois ambientes de testes de custos — Cena 1 (Nuvem paga) vs Cena 2 (Kronk local)",
            "Visualização direta — acompanhe os contadores integrados de tokens e custos na UI"
        ],
        "notes": "Lembram do gancho do começo — 'tem um jeito ingênuo e um "
                 "jeito inteligente'? Vou rodar exatamente a mesma tarefa "
                 "duas vezes, mudando só o MCP ativo e o modelo por trás."
    },
    {
        "title": "Slide 21 — Apresentando a Cena 1",
        "bullets": [
            "Backend do Provedor — API de nuvem comercial do parceiro integrada com créditos de uso",
            "Servidor MCP ativo — vault_dumb escutando na porta local :9001",
            "Preparação dos contadores — zerados para garantir transparência na medição dos gastos"
        ],
        "notes": "Cena 1 representa o jeito mais comum hoje: modelo na "
                 "nuvem, pago por token, conectado a um MCP que não "
                 "filtra nada. Olhem os números antes de eu começar: "
                 "zero."
    },
    {
        "title": "Slide 22 — Cena 1 ao vivo  [DEMO]",
        "bullets": [
            "O prompt de comando — 'Leia a TASK-001 no vault do Obsidian e resolva o problema de concorrência'",
            "Decisão autônoma — o agente escolhe chamar read_vault (única ferramenta disponível)",
            "Volume de processamento — observe a subida massiva do contador de tokens e custo disparando na tela"
        ],
        "notes": "Reparem: o Cline está decidindo chamar read_vault — só "
                 "tem essa ferramenta disponível. Contingência: se a "
                 "geração de nuvem travar num tool-call depois, o ponto "
                 "já foi feito no instante em que read_vault devolve a "
                 "resposta — pode pausar aí e seguir pra Cena 2."
    },
    {
        "title": "Slide 23 — Métricas de Consumo: Cena 1",
        "bullets": [
            "Controle de custos reais — exibição direta do consumo financeiro e de volume medidos no terminal",
            "O impacto do desperdício — volume expressivo de tokens consumidos apenas na coleta do contexto",
            "Análise do gargalo — a simplicidade de código do dumb mcp penaliza financeiramente o negócio"
        ],
        "notes": "Aqui está o resultado da Cena 1: [ler os números reais "
                 "capturados]. Token e dinheiro de verdade, gastos só "
                 "porque a ferramenta de busca não soube filtrar nada. "
                 "Guardem esse número."
    },
    {
        "title": "Slide 24 — Apresentando a Cena 2",
        "bullets": [
            "Infraestrutura de inferência — Kronk rodando localmente na máquina do desenvolvedor (:11435)",
            "Servidor MCP ativo — vault_smart habilitado de forma exclusiva na porta local :9002",
            "Simetria de teste — mesmo prompt textual colado exatamente de forma idêntica"
        ],
        "notes": "Troquei o provider do Cline pra apontar pro Kronk, "
                 "rodando aqui, agora, sem internet envolvida. E troquei "
                 "o MCP ativo pra vault_smart."
    },
    {
        "title": "Slide 25 — Cena 2 ao vivo  [DEMO]",
        "bullets": [
            "Início do fluxo — envio do prompt idêntico sem acesso de rede externa",
            "Chamada qualificada — agente Cline decide consultar a ferramenta search_vault",
            "Processamento cirúrgico — o modelo local recebe apenas os 3 documentos com semântica compatível",
            "Custo financeiro absoluto — zero dólares e zero centavos consumidos rodando de forma local-first"
        ],
        "notes": "Reparem que a ferramenta disponível agora é outra: "
                 "search_vault. Contingência: o modelo local com "
                 "raciocínio pode demorar mais que a nuvem pra terminar "
                 "o fix completo — o ponto da demo é o contraste de "
                 "tokens/custo na chamada da ferramenta, que aparece "
                 "rápido."
    },
    {
        "title": "Slide 26 — Comparativo Consolidado: Cena 1 vs Cena 2",
        "table": [
            ["Comparativo", "Cena 1 (Nuvem + Dumb MCP)", "Cena 2 (Kronk + Smart MCP)"],
            ["Tokens Consumidos", "Massivo (~47.000 tokens)", "Otimizado (~4.000 tokens)"],
            ["Custo Financeiro", "Cobrança real em dólares", "$0.00 (Processamento local)"],
            ["Privacidade de Dados", "Código enviado para servidores externos", "100% contido localmente nesta máquina"],
            ["Necessidade de Rede", "Acesso à internet obrigatório", "Totalmente offline (Wi-Fi desativado)"],
            ["Qualidade de Resposta", "Modelo pode se dispersar no backlog", "Foco absoluto no contexto cirúrgico"]
        ],
        "notes": "Essa é a imagem que eu queria que vocês levassem pra "
                 "casa. Mesma tarefa, mesmo agente. A diferença está em "
                 "como a informação chega até o modelo, e onde o modelo "
                 "roda."
    },
    {
        "title": "Slide 27 — Fechamento",
        "gopher": True,
        "bullets": [
            "Go como infraestrutura — Go é de longe a linguagem mais recomendada para orquestrar dados de IA locais",
            "Independência real — construímos pipelines, embeddings e bancos vetoriais sem recorrer a pilhas externas",
            "Código livre — todos os testes, MCPs e este deck estão disponíveis em repositório público",
            "Repositório da Demo — github.com/jeffersonferreira/kronk-tech-demo-golangsp"
        ],
        "notes": "Tudo que vocês viram hoje — servidor MCP, motor de "
                 "inferência, busca vetorial — é Go, de ponta a ponta. "
                 "Não precisamos de Python, não precisamos de nuvem. "
                 "Obrigado!"
    }
]


def set_slide_background(slide, color):
    """Sets a solid background color to the slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def draw_header_footer(slide, title, category, slide_num, total_slides):
    """Draws a professional, high-end technical header and footer."""
    # 1. Slide Category Tag (top-left) - Electric Cyan for maximum contrast on deep blue
    tag_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.25))
    tf_tag = tag_box.text_frame
    tf_tag.word_wrap = True
    tf_tag.margin_left = tf_tag.margin_top = tf_tag.margin_right = tf_tag.margin_bottom = 0
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = category.upper()
    p_tag.font.name = "Segoe UI"
    p_tag.font.size = Pt(9)
    p_tag.font.bold = True
    p_tag.font.color.rgb = RGBColor(0x00, 0xD2, 0xFF) # Electric Gopher Cyan

    # 2. Main Slide Title - Pure White for maximum legibility on deep blue background
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.55), Inches(12.0), Inches(0.5))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) # High-contrast Pure White

    # 3. Horizontal Separator Line (top) - Subtle deep navy divider
    sep_color = RGBColor(0x1E, 0x3A, 0x5F)
    sep_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.15), Inches(12.133), Inches(0.01))
    sep_line.fill.solid()
    sep_line.fill.fore_color.rgb = sep_color
    sep_line.line.color.rgb = sep_color

    # 4. Horizontal Footer Separator Line - Subtle deep navy divider
    footer_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.85), Inches(12.133), Inches(0.01))
    footer_line.fill.solid()
    footer_line.fill.fore_color.rgb = sep_color
    footer_line.line.color.rgb = sep_color

    # 5. Footer Left Text - Legible silver-gray on deep blue background
    f_left_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.95), Inches(6.0), Inches(0.3))
    tf_f_left = f_left_box.text_frame
    tf_f_left.margin_left = tf_f_left.margin_top = tf_f_left.margin_right = tf_f_left.margin_bottom = 0
    p_f_left = tf_f_left.paragraphs[0]
    p_f_left.text = "GolangSP 2026  |  IA sem Nuvem, em Golang"
    p_f_left.font.name = "Segoe UI"
    p_f_left.font.size = Pt(8.5)
    p_f_left.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8) # Slate-400

    # 6. Footer Right Page Number - Legible silver-gray on deep blue background
    f_right_box = slide.shapes.add_textbox(Inches(10.0), Inches(6.95), Inches(2.733), Inches(0.3))
    tf_f_right = f_right_box.text_frame
    tf_f_right.margin_left = tf_f_right.margin_top = tf_f_right.margin_right = tf_f_right.margin_bottom = 0
    p_f_right = tf_f_right.paragraphs[0]
    p_f_right.alignment = PP_ALIGN.RIGHT
    p_f_right.text = f"{slide_num} / {total_slides}"
    p_f_right.font.name = "Segoe UI"
    p_f_right.font.size = Pt(8.5)
    p_f_right.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8) # Slate-400


def format_bullet_text(paragraph, text):
    """Applies smart typographic formatting to highlight key terms before separators."""
    paragraph.font.name = "Segoe UI"
    paragraph.font.size = Pt(13)
    paragraph.line_spacing = 1.25
    paragraph.space_after = Pt(10)
    
    # Check for visual separators to split and highlight key terms
    separators = [" — ", ": ", " · ", " -> ", " → "]
    found_sep = None
    for sep in separators:
        if sep in text:
            found_sep = sep
            break
            
    if found_sep:
        parts = text.split(found_sep, 1)
        prefix = parts[0] + found_sep
        suffix = parts[1]
        
        # Highlighted Prefix (Royal Blue & Bold) - Beautiful contrast on White cards
        run_pre = paragraph.add_run()
        run_pre.text = prefix
        run_pre.font.bold = True
        run_pre.font.color.rgb = ACCENT_BLUE
        
        # Plain Body Suffix (Muted Charcoal Slate-600)
        run_suf = paragraph.add_run()
        run_suf.text = suffix
        run_suf.font.bold = False
        run_suf.font.color.rgb = TEXT_MUTED
    else:
        # Standard bullet style
        run = paragraph.add_run()
        run.text = text
        run.font.bold = False
        run.font.color.rgb = TEXT_MAIN


def add_modern_bullets(tf, bullets):
    """Populates bullet list using premium formatting styles."""
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    # Process first bullet
    p_first = tf.paragraphs[0]
    format_bullet_text(p_first, bullets[0])
    
    # Process remaining bullets
    for b in bullets[1:]:
        p = tf.add_paragraph()
        p.space_before = Pt(8)
        format_bullet_text(p, b)


def draw_terminal_window(slide, left, top, width, height, filename, code):
    """Draws a premium VS Code-like monospaced terminal window container."""
    # 1. Main terminal body card
    terminal_body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    terminal_body.fill.solid()
    terminal_body.fill.fore_color.rgb = CODE_BG
    terminal_body.line.color.rgb = CODE_BORDER
    terminal_body.line.width = Pt(1.5)

    # 2. Terminal title header bar
    top_bar_height = Inches(0.35)
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, top_bar_height)
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = CARD_BORDER
    top_bar.line.color.rgb = CARD_BORDER

    # 3. Three window control button circles (macOS style)
    btn_size = Inches(0.09)
    btn_y = top + Inches(0.13)
    colors = [
        RGBColor(0xFF, 0x5F, 0x56), # Red
        RGBColor(0xFF, 0xBD, 0x2E), # Yellow
        RGBColor(0x27, 0xC9, 0x3F)  # Green
    ]
    for i, color in enumerate(colors):
        btn_x = left + Inches(0.15) + (i * Inches(0.14))
        btn = slide.shapes.add_shape(MSO_SHAPE.OVAL, btn_x, btn_y, btn_size, btn_size)
        btn.fill.solid()
        btn.fill.fore_color.rgb = color
        btn.line.color.rgb = color

    # 4. Monospace File Tab Name (center-aligned in top bar)
    tab_box = slide.shapes.add_textbox(left + Inches(1.0), top + Inches(0.04), width - Inches(2.0), Inches(0.25))
    tf_tab = tab_box.text_frame
    tf_tab.margin_left = tf_tab.margin_top = tf_tab.margin_right = tf_tab.margin_bottom = 0
    p_tab = tf_tab.paragraphs[0]
    p_tab.alignment = PP_ALIGN.CENTER
    p_tab.text = filename
    p_tab.font.name = "Consolas"
    p_tab.font.size = Pt(9.5)
    p_tab.font.bold = True
    p_tab.font.color.rgb = TEXT_MUTED

    # 5. Monospace Code Text Box
    code_box = slide.shapes.add_textbox(left + Inches(0.18), top + top_bar_height + Inches(0.15), width - Inches(0.36), height - top_bar_height - Inches(0.25))
    tf_code = code_box.text_frame
    tf_code.word_wrap = True
    tf_code.margin_left = tf_code.margin_top = tf_code.margin_right = tf_code.margin_bottom = 0
    
    lines = code.split("\n")
    p_first_line = tf_code.paragraphs[0]
    p_first_line.text = lines[0]
    p_first_line.font.name = "Consolas"
    p_first_line.font.size = Pt(9.5)
    p_first_line.font.color.rgb = TEXT_MAIN
    
    for line in lines[1:]:
        p = tf_code.add_paragraph()
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(9.5)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(2)


def draw_styled_table(slide, rows, left, top, width, height, side_by_side=False):
    """Creates a beautifully styled modern UI table with corporate SaaS dashboard aesthetic on white cards."""
    n_rows, n_cols = len(rows), len(rows[0])
    gtable = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    
    font_size = Pt(10) if side_by_side else Pt(11)
    header_font_size = Pt(11) if side_by_side else Pt(12)
    
    for r in range(n_rows):
        for c in range(n_cols):
            cell = gtable.cell(r, c)
            
            # 1. Custom Background Styling for Cell - Stripe & SaaS aesthetic
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = RGBColor(0x0A, 0x25, 0x40) # Stripe Deep Navy header
            elif r % 2 == 1:
                cell.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC) # Light gray zebra alternate row
            else:
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) # Pure White
                
            # 2. Text Frame Properties
            tf_cell = cell.text_frame
            tf_cell.word_wrap = True
            
            p = tf_cell.paragraphs[0]
            p.text = rows[r][c]
            p.alignment = PP_ALIGN.LEFT
            
            # Align second, third, fourth columns containing short stats to center
            if c > 0 and n_cols > 2:
                p.alignment = PP_ALIGN.CENTER
                
            # 3. Apply Professional Typography
            for run in p.runs:
                run.font.name = "Segoe UI"
                if r == 0:
                    run.font.bold = True
                    run.font.size = header_font_size
                    run.font.color.rgb = RGBColor(0x00, 0xD2, 0xFF)     # Electric Gopher Cyan text on dark header
                else:
                    run.font.bold = False
                    run.font.size = font_size
                    run.font.color.rgb = TEXT_MAIN                      # Deep Charcoal Slate-900 text on white/gray cells


def draw_bio_content(slide, name, role, companies, content_top, content_height):
    """Draws a premium speaker profile dashboard with a photo frame and interactive-looking company timeline cards."""
    # 1. Left Side: Photo Frame Container (Sleek dark panel containing centered square photo and social links)
    left_frame = Inches(0.6)
    width_frame = Inches(4.5)
    
    # Outer Frame Card
    photo_card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_frame, content_top, width_frame, content_height)
    photo_card.fill.solid()
    photo_card.fill.fore_color.rgb = CODE_BG
    photo_card.line.color.rgb = CARD_BORDER
    photo_card.line.width = Pt(1.5)
    
    # Center-aligned 1:1 square photo mount (looks like a photo slide)
    photo_size = Inches(2.8)
    photo_mount = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        left_frame + (width_frame - photo_size) / 2, 
        content_top + Inches(0.4), 
        photo_size, 
        photo_size
    )
    photo_mount.fill.solid()
    photo_mount.fill.fore_color.rgb = CARD_BG # Crisp pure White background for the photo frame
    photo_mount.line.color.rgb = CARD_BORDER
    photo_mount.line.width = Pt(1.5)
    
    import os
    me_photo_path = "slides/assets/me.jpeg"
    if os.path.exists(me_photo_path):
        # Physically insert the real photo centered inside the white mount frame
        photo_border = Inches(0.12)
        slide.shapes.add_picture(
            me_photo_path,
            left_frame + (width_frame - photo_size) / 2 + photo_border,
            content_top + Inches(0.4) + photo_border,
            width=photo_size - (2 * photo_border),
            height=photo_size - (2 * photo_border)
        )
    else:
        # Photo placeholder text instructions inside the square if file is missing
        txt_box = slide.shapes.add_textbox(
            left_frame + (width_frame - photo_size) / 2 + Inches(0.15), 
            content_top + Inches(1.1), 
            photo_size - Inches(0.3), 
            Inches(1.2)
        )
        tf_txt = txt_box.text_frame
        tf_txt.word_wrap = True
        p1 = tf_txt.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        p1.text = "[ SUA FOTO 1:1 ]"
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = ACCENT_BLUE
        
        p2 = tf_txt.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(6)
        p2.text = "Arraste e solte uma imagem quadrada aqui."
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = TEXT_MUTED

    # Speaker Social Info Badge (underneath the square photo frame)
    info_box = slide.shapes.add_textbox(left_frame + Inches(0.2), content_top + Inches(3.45), width_frame - Inches(0.4), Inches(1.5))
    tf_info = info_box.text_frame
    tf_info.word_wrap = True
    tf_info.margin_left = tf_info.margin_top = tf_info.margin_right = tf_info.margin_bottom = 0
    
    p_info_name = tf_info.paragraphs[0]
    p_info_name.alignment = PP_ALIGN.CENTER
    p_info_name.text = name
    p_info_name.font.name = "Segoe UI"
    p_info_name.font.size = Pt(13)
    p_info_name.font.bold = True
    p_info_name.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) # High-contrast White name inside dark card
    
    p_git = tf_info.add_paragraph()
    p_git.alignment = PP_ALIGN.CENTER
    p_git.space_before = Pt(5)
    p_git.text = "GitHub: github.com/jeffersonferreira"
    p_git.font.name = "Consolas"
    p_git.font.size = Pt(8.5)
    p_git.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8) # Slate-400
    
    p_ln = tf_info.add_paragraph()
    p_ln.alignment = PP_ALIGN.CENTER
    p_ln.space_before = Pt(3)
    p_ln.text = "LinkedIn: linkedin.com/in/jefferson-f"
    p_ln.font.name = "Consolas"
    p_ln.font.size = Pt(8.5)
    p_ln.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8) # Slate-400

    # 2. Right Side: Bio and Companies
    left_col = Inches(5.4)
    width_col = Inches(7.333)
    
    # Name Text Frame
    name_box = slide.shapes.add_textbox(left_col, content_top, width_col, Inches(1.0))
    tf_name = name_box.text_frame
    tf_name.word_wrap = True
    tf_name.margin_left = tf_name.margin_top = tf_name.margin_right = tf_name.margin_bottom = 0
    
    # Speaker Name
    p_name = tf_name.paragraphs[0]
    p_name.text = name
    p_name.font.name = "Segoe UI"
    p_name.font.size = Pt(30)
    p_name.font.bold = True
    p_name.font.color.rgb = TEXT_MAIN
    
    # Speaker Role
    p_role = tf_name.add_paragraph()
    p_role.text = role
    p_role.font.name = "Segoe UI"
    p_role.font.size = Pt(14)
    p_role.font.bold = True
    p_role.font.color.rgb = ACCENT_BLUE
    p_role.space_before = Pt(4)
    
    # Timeline Header
    timeline_top = content_top + Inches(1.1)
    hdr_box = slide.shapes.add_textbox(left_col, timeline_top, width_col, Inches(0.3))
    tf_hdr = hdr_box.text_frame
    tf_hdr.margin_left = tf_hdr.margin_top = tf_hdr.margin_right = tf_hdr.margin_bottom = 0
    p_hdr = tf_hdr.paragraphs[0]
    p_hdr.text = "TRAJETÓRIA PROFISSIONAL RELEVANTE"
    p_hdr.font.name = "Segoe UI"
    p_hdr.font.size = Pt(9.5)
    p_hdr.font.bold = True
    p_hdr.font.color.rgb = TEXT_MUTED
    
    # Draw Company Cards
    card_start_y = timeline_top + Inches(0.35)
    card_height = Inches(1.0)
    card_spacing = Inches(0.18)
    
    for i, comp in enumerate(companies):
        card_y = card_start_y + (i * (card_height + card_spacing))
        
        # Company Card container
        comp_card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_col, card_y, width_col, card_height)
        comp_card.fill.solid()
        comp_card.fill.fore_color.rgb = CARD_BG
        comp_card.line.color.rgb = CARD_BORDER
        comp_card.line.width = Pt(1.0)
        
        # Glow strip on the left of each card
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_col, card_y, Inches(0.08), card_height)
        strip.fill.solid()
        strip.fill.fore_color.rgb = ACCENT_BLUE if "1Global" in comp["name"] else CARD_BORDER
        strip.line.color.rgb = ACCENT_BLUE if "1Global" in comp["name"] else CARD_BORDER
        
        # Text inside company card
        comp_box = slide.shapes.add_textbox(left_col + Inches(0.25), card_y + Inches(0.12), width_col - Inches(0.4), card_height - Inches(0.2))
        tf_comp = comp_box.text_frame
        tf_comp.word_wrap = True
        tf_comp.margin_left = tf_comp.margin_top = tf_comp.margin_right = tf_comp.margin_bottom = 0
        
        # Company Name & Status
        p_comp_name = tf_comp.paragraphs[0]
        p_comp_name.text = comp["name"]
        p_comp_name.font.name = "Segoe UI"
        p_comp_name.font.size = Pt(12)
        p_comp_name.font.bold = True
        p_comp_name.font.color.rgb = ACCENT_BLUE if "1Global" in comp["name"] else TEXT_MAIN
        
        # Company Description / Role
        p_comp_desc = tf_comp.add_paragraph()
        p_comp_desc.space_before = Pt(4)
        p_comp_desc.text = comp["desc"]
        p_comp_desc.font.name = "Segoe UI"
        p_comp_desc.font.size = Pt(10.5)
        p_comp_desc.font.color.rgb = TEXT_MUTED if "1Global" not in comp["name"] else TEXT_MAIN


def build_title_slide(slide, title, subtitle):
    """Renders a pixel-perfect, premium tech-themed title slide."""
    set_slide_background(slide, BG_COLOR)

    # 1. Glassmorphic main container panel
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(1.3), Inches(11.633), Inches(4.9))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = CARD_BORDER
    card.line.width = Pt(1.5)

    # 2. Go Cyan glowing border stripe on left card boundary
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(1.3), Inches(0.15), Inches(4.9))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT_BLUE
    stripe.line.color.rgb = ACCENT_BLUE

    # 3. Event Tag (top-left inside card)
    tag_box = slide.shapes.add_textbox(Inches(1.4), Inches(1.8), Inches(7.0), Inches(0.3))
    tf_tag = tag_box.text_frame
    tf_tag.margin_left = tf_tag.margin_top = tf_tag.margin_right = tf_tag.margin_bottom = 0
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = "GOLANGSP SPECIAL PRESENTATION  ·  PALCO PRINCIPAL"
    p_tag.font.name = "Segoe UI"
    p_tag.font.size = Pt(10.5)
    p_tag.font.bold = True
    p_tag.font.color.rgb = ACCENT_BLUE

    # 4. Main Talk Title (large bold)
    title_box = slide.shapes.add_textbox(Inches(1.4), Inches(2.2), Inches(7.2), Inches(1.2))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(44)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_MAIN

    # 5. Secondary Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1.4), Inches(3.5), Inches(7.2), Inches(0.8))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = subtitle
    p_sub.font.name = "Segoe UI"
    p_sub.font.size = Pt(17)
    p_sub.font.color.rgb = TEXT_MUTED

    # 6. Speaker details (bottom-left inside card)
    spk_box = slide.shapes.add_textbox(Inches(1.4), Inches(4.75), Inches(7.0), Inches(0.6))
    tf_spk = spk_box.text_frame
    tf_spk.margin_left = tf_spk.margin_top = tf_spk.margin_right = tf_spk.margin_bottom = 0
    p_spk = tf_spk.paragraphs[0]
    p_spk.text = "Apresentado por Jefferson Ferreira"
    p_spk.font.name = "Segoe UI"
    p_spk.font.size = Pt(14)
    p_spk.font.bold = True
    p_spk.font.color.rgb = TEXT_MAIN
    
    p_spk_sub = tf_spk.add_paragraph()
    p_spk_sub.text = "IA sem Nuvem, Local-First, Alta Performance em Go"
    p_spk_sub.font.name = "Segoe UI"
    p_spk_sub.font.size = Pt(11)
    p_spk_sub.font.color.rgb = ACCENT_BLUE

    # 7. Go Gopher Image (right alignment inside card)
    slide.shapes.add_picture(GOPHER_PATH, Inches(9.1), Inches(1.85), width=Inches(2.8))


def build_section_divider_slide(slide, number_str, title_str, subtitle_str):
    """Creates a high-impact Section Divider slide to transition presentation chapters."""
    set_slide_background(slide, BG_COLOR)

    # 1. Subtle Background Chapter Watermark - Deep blue blending softly with BG_COLOR
    wm_box = slide.shapes.add_textbox(Inches(6.4), Inches(0.8), Inches(6.5), Inches(5.5))
    tf_wm = wm_box.text_frame
    tf_wm.word_wrap = True
    tf_wm.margin_left = tf_wm.margin_top = tf_wm.margin_right = tf_wm.margin_bottom = 0
    p_wm = tf_wm.paragraphs[0]
    p_wm.alignment = PP_ALIGN.RIGHT
    p_wm.text = number_str
    p_wm.font.name = "Segoe UI"
    p_wm.font.size = Pt(180)
    p_wm.font.bold = True
    p_wm.font.color.rgb = RGBColor(0x13, 0x34, 0x59) # Subtle blending blue

    # 2. Main foreground section title content
    content_box = slide.shapes.add_textbox(Inches(1.1), Inches(2.2), Inches(7.5), Inches(4.0))
    tf_content = content_box.text_frame
    tf_content.word_wrap = True
    tf_content.margin_left = tf_content.margin_top = tf_content.margin_right = tf_content.margin_bottom = 0
    
    # Category Header - Electric Gopher Cyan
    p_cat = tf_content.paragraphs[0]
    p_cat.text = f"CONTEÚDO DA APRESENTAÇÃO  ·  CAPÍTULO {number_str}"
    p_cat.font.name = "Segoe UI"
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = RGBColor(0x00, 0xD2, 0xFF)
    p_cat.space_after = Pt(12)
    
    # Title - Pure White
    p_title = tf_content.add_paragraph()
    p_title.text = title_str
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p_title.space_after = Pt(15)
    
    # Decorative line indicator - Electric Gopher Cyan
    dec_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.1), Inches(4.35), Inches(1.8), Inches(0.04))
    dec_line.fill.solid()
    dec_line.fill.fore_color.rgb = RGBColor(0x00, 0xD2, 0xFF)
    dec_line.line.color.rgb = RGBColor(0x00, 0xD2, 0xFF)
    
    # Subtitle / description - Cool slate-silver
    p_sub = tf_content.add_paragraph()
    p_sub.space_before = Pt(15)
    p_sub.text = subtitle_str
    p_sub.font.name = "Segoe UI"
    p_sub.font.size = Pt(15)
    p_sub.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)


def get_slide_section(title, index):
    """Maps content slides to structural chapters."""
    if index <= 13:
        return "intro"
    elif index <= 30:
        return "kronk"
    elif index <= 41:
        return "code"
    elif index <= 48:
        return "demo"
    else:
        return "outro"


def get_slide_category_name(section):
    """Translates structural code to visible headers on content slides."""
    mapping = {
        "intro": "Introdução & Contexto",
        "kronk": "Mergulho no Kronk SDK",
        "code": "Por Dentro do Código MCP",
        "demo": "A Prova Real - Demos",
        "outro": "Fechamento da Talk"
    }
    return mapping.get(section, "Apresentação")


def build():
    """Builds the entire custom PPTX deck."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6] # Blank Layout

    # Filter out initial main title spec to find out correct content slides count
    total_content_slides = sum(1 for spec in SLIDES if spec.get("kind") != "title")
    
    current_section = None
    slide_counter = 0

    for idx, spec in enumerate(SLIDES):
        # 1. Main Cover Slide
        if spec.get("kind") == "title":
            slide = prs.slides.add_slide(blank_layout)
            build_title_slide(slide, spec["title"], spec["subtitle"])
            
            # Speaker notes for title slide
            notes = spec.get("notes")
            if notes:
                slide.notes_slide.notes_text_frame.text = notes
            continue

        # 2. Resolve Chapter / Section transitions
        slide_counter += 1
        sec = get_slide_section(spec["title"], idx)
        
        if sec != current_section:
            # Transition occurred - Insert Section Divider first!
            current_section = sec
            meta = SECTIONS_META[current_section]
            div_slide = prs.slides.add_slide(blank_layout)
            build_section_divider_slide(div_slide, meta["num"], meta["title"], meta["sub"])
        
        # 3. Create actual content slide
        slide = prs.slides.add_slide(blank_layout)
        set_slide_background(slide, BG_COLOR)

        # Sanitize drafts title prefixes e.g. "Slide 1 — O que é o Kronk?" -> "O que é o Kronk?"
        clean_title = spec["title"]
        clean_title = re.sub(r"^Slide\s+\d+\s+[-—]\s+", "", clean_title)

        category_name = get_slide_category_name(current_section)
        if spec.get("kind") == "bio":
            category_name = "Palestrante"
            
        draw_header_footer(slide, clean_title, category_name, slide_counter, total_content_slides)

        has_code = "code" in spec
        has_table = "table" in spec
        bullets = spec.get("bullets")

        # Layout metrics inside Content Area
        content_top = Inches(1.4)
        content_height = Inches(5.3)

        # Case BIO: Custom Speaker Profile Layout
        if spec.get("kind") == "bio":
            draw_bio_content(slide, spec["title"], spec["role"], spec["companies"], content_top, content_height)

        # Case A: Side-by-Side Bullets on Left + Code/Table on Right
        elif bullets and (has_code or has_table):
            # Left Card Container for bullets
            left_card_width = Inches(5.8)
            l_card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), content_top, left_card_width, content_height)
            l_card.fill.solid()
            l_card.fill.fore_color.rgb = CARD_BG
            l_card.line.color.rgb = CARD_BORDER
            l_card.line.width = Pt(1.5)

            # Left Card text box padding
            bullets_box = slide.shapes.add_textbox(Inches(0.95), content_top + Inches(0.3), left_card_width - Inches(0.65), content_height - Inches(0.6))
            add_modern_bullets(bullets_box.text_frame, bullets)

            # Right Side Components
            right_left = Inches(6.833)
            right_width = Inches(5.9)
            
            if has_code:
                filename = spec.get("code_lang", "terminal").upper()
                if filename == "BASH":
                    filename = "terminal"
                elif filename == "JSON":
                    filename = "response.json"
                elif filename == "YAML":
                    filename = "model_config.yaml"
                elif filename == "GO":
                    # Simple heuristic to name custom tabs
                    if "LoadAll" in spec["code"]:
                        filename = "internal/vault/reader.go"
                    elif "readVaultHandler" in spec["code"]:
                        filename = "cmd/dumb/main.go"
                    else:
                        filename = "cmd/smart/main.go"
                elif filename == "SQL":
                    filename = "search_vault.sql"
                    
                draw_terminal_window(slide, right_left, content_top, right_width, content_height, filename, spec["code"])
            elif has_table:
                draw_styled_table(slide, spec["table"], right_left, content_top + Inches(0.2), right_width, Inches(0.4 * len(spec["table"])), side_by_side=True)

        # Case B: Full-Width Code/Table with NO bullets
        elif (has_code or has_table) and not bullets:
            center_left = Inches(0.6)
            center_width = Inches(12.133)
            
            if has_code:
                filename = spec.get("code_lang", "file").upper()
                draw_terminal_window(slide, center_left, content_top, center_width, content_height, filename, spec["code"])
            elif has_table:
                draw_styled_table(slide, spec["table"], center_left, content_top + Inches(0.5), center_width, Inches(0.5 * len(spec["table"])), side_by_side=False)

        # Case C: Full-Width Bullets Only
        elif bullets and not (has_code or has_table):
            center_left = Inches(0.6)
            center_width = Inches(12.133)
            
            # Massive card container for bullets
            main_card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, center_left, content_top, center_width, content_height)
            main_card.fill.solid()
            main_card.fill.fore_color.rgb = CARD_BG
            main_card.line.color.rgb = CARD_BORDER
            main_card.line.width = Pt(1.5)

            bullets_box = slide.shapes.add_textbox(center_left + Inches(0.45), content_top + Inches(0.4), center_width - Inches(0.9), content_height - Inches(0.8))
            add_modern_bullets(bullets_box.text_frame, bullets)

        # 4. Handle Gopher Asset on the closing slide
        if spec.get("gopher"):
            slide.shapes.add_picture(GOPHER_PATH, Inches(9.8), Inches(2.2), width=Inches(2.6))

        # 5. Populate Spoken Speaker Notes
        notes = spec.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(OUT_PATH)
    print(f"Wrote {len(prs.slides._sldIdLst)} slides to {OUT_PATH}")


if __name__ == "__main__":
    build()
